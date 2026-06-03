"""
Парсеры и генераторы файлов.

Чтение:
- .txt/.md/.csv/.json/.log/.xml/.yaml/.py/.js/.html/.css — read text
- .pdf  → pypdf
- .docx → python-docx
- .xlsx → openpyxl

Запись:
- .txt/.md/.log     — utf-8
- .csv              — passthrough (Claude передаёт уже валидный csv)
- .json             — валидируем + pretty
- .docx             — markdown-like → docx (заголовки, списки, абзацы)
- .xlsx             — csv-строка → xlsx таблица

Изображения сюда не входят — они идут vision-каналом (см. agent.py).
"""

from __future__ import annotations

import csv
import io
import json
import mimetypes
from pathlib import Path

from bot.config import MAX_EXTRACTED_CHARS
from bot.logger import get_logger


logger = get_logger("files")


TEXT_EXT = {
    ".txt", ".md", ".csv", ".json", ".log", ".xml",
    ".yaml", ".yml", ".py", ".js", ".html", ".css",
}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
PARSE_EXT = {".pdf", ".docx", ".xlsx"}
ALLOWED_INPUT_EXT = TEXT_EXT | IMAGE_EXT | PARSE_EXT
GENERATE_EXT = {".txt", ".md", ".log", ".csv", ".json", ".docx", ".xlsx", ".pptx"}


def detect_ext(name: str) -> str:
    return Path(name).suffix.lower()


def is_allowed_input(name: str) -> bool:
    return detect_ext(name) in ALLOWED_INPUT_EXT


def is_image(name: str) -> bool:
    return detect_ext(name) in IMAGE_EXT


def guess_mime(name: str) -> str:
    return mimetypes.guess_type(name)[0] or "application/octet-stream"


def _truncate(s: str) -> str:
    if len(s) > MAX_EXTRACTED_CHARS:
        return s[:MAX_EXTRACTED_CHARS] + f"\n\n[...обрезано, всего {len(s)} символов]"
    return s


# ---------------- READ ----------------

def extract_text(path: str | Path) -> str:
    """Извлечь текст из файла. Никогда не выбрасывает — возвращает текст ошибки."""
    p = Path(path)
    if not p.exists():
        return f"[Файл не найден: {p.name}]"
    ext = p.suffix.lower()

    if ext in TEXT_EXT:
        try:
            return _truncate(p.read_text(encoding="utf-8"))
        except UnicodeDecodeError:
            try:
                return _truncate(p.read_text(encoding="cp1251", errors="replace"))
            except OSError as e:
                return f"[Ошибка чтения: {e}]"
        except OSError as e:
            return f"[Ошибка чтения: {e}]"

    if ext == ".pdf":
        try:
            from pypdf import PdfReader  # type: ignore
        except ImportError:
            return "[Не установлен pypdf]"
        try:
            reader = PdfReader(str(p))
            chunks: list[str] = []
            for i, page in enumerate(reader.pages):
                try:
                    chunks.append(page.extract_text() or "")
                except Exception as e:
                    chunks.append(f"[стр.{i + 1}: ошибка {e.__class__.__name__}]")
            return _truncate("\n".join(chunks).strip() or "[PDF пустой или текст не извлекается]")
        except Exception as e:
            return f"[Не удалось прочитать PDF: {e}]"

    if ext == ".docx":
        try:
            from docx import Document  # type: ignore
        except ImportError:
            return "[Не установлен python-docx]"
        try:
            doc = Document(str(p))
            chunks = [para.text for para in doc.paragraphs if para.text]
            for table in doc.tables:
                for row in table.rows:
                    chunks.append(" | ".join(c.text for c in row.cells))
            return _truncate("\n".join(chunks).strip() or "[DOCX пустой]")
        except Exception as e:
            return f"[Не удалось прочитать DOCX: {e}]"

    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook  # type: ignore
        except ImportError:
            return "[Не установлен openpyxl]"
        try:
            wb = load_workbook(str(p), data_only=True, read_only=True)
            chunks: list[str] = []
            for ws in wb.worksheets:
                chunks.append(f"--- Лист: {ws.title} ---")
                for row in ws.iter_rows(values_only=True):
                    chunks.append(" | ".join("" if c is None else str(c) for c in row))
            wb.close()
            return _truncate("\n".join(chunks).strip() or "[XLSX пустой]")
        except Exception as e:
            return f"[Не удалось прочитать XLSX: {e}]"

    return f"[Формат {ext or 'без расширения'} не поддерживается для чтения]"


# ---------------- WRITE ----------------

def generate_file(name: str, content: str) -> tuple[bytes, str]:
    """
    Сгенерировать файл по имени и содержимому.
    Возвращает (bytes, mime). Бросает ValueError при неподдерживаемом формате.
    """
    ext = detect_ext(name)
    if ext not in GENERATE_EXT:
        raise ValueError(
            f"Формат {ext or 'без расширения'} нельзя создать. "
            f"Доступно: {', '.join(sorted(GENERATE_EXT))}"
        )

    if ext in {".txt", ".md", ".log"}:
        return content.encode("utf-8"), "text/plain; charset=utf-8"

    if ext == ".csv":
        return content.encode("utf-8"), "text/csv; charset=utf-8"

    if ext == ".json":
        try:
            parsed = json.loads(content)
            pretty = json.dumps(parsed, ensure_ascii=False, indent=2)
            return pretty.encode("utf-8"), "application/json"
        except json.JSONDecodeError:
            return content.encode("utf-8"), "application/json"

    if ext == ".docx":
        from docx import Document  # type: ignore

        doc = Document()
        for raw_line in content.split("\n"):
            line = raw_line.rstrip()
            if not line:
                doc.add_paragraph("")
                continue
            if line.startswith("# "):
                doc.add_heading(line[2:].strip(), level=1)
            elif line.startswith("## "):
                doc.add_heading(line[3:].strip(), level=2)
            elif line.startswith("### "):
                doc.add_heading(line[4:].strip(), level=3)
            elif line.startswith(("- ", "* ")):
                doc.add_paragraph(line[2:].strip(), style="List Bullet")
            elif line[:2].rstrip(".").isdigit() and line[1:3].startswith(". "):
                doc.add_paragraph(line.split(". ", 1)[1], style="List Number")
            else:
                doc.add_paragraph(line)
        buf = io.BytesIO()
        doc.save(buf)
        return (
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

    if ext == ".xlsx":
        from openpyxl import Workbook  # type: ignore

        wb = Workbook()
        ws = wb.active
        reader = csv.reader(io.StringIO(content))
        for row in reader:
            ws.append(row)
        buf = io.BytesIO()
        wb.save(buf)
        return (
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    if ext == ".pptx":
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Pt  # type: ignore
        except ImportError:
            raise ValueError("Не установлен python-pptx. Запусти: pip install python-pptx")

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]   # blank
        title_layout = prs.slide_layouts[0]   # title slide
        content_layout = prs.slide_layouts[1] # title + content

        slides_raw = content.split("\n---\n")
        for idx, slide_text in enumerate(slides_raw):
            lines = [l.rstrip() for l in slide_text.strip().split("\n") if l.strip()]
            if not lines:
                continue

            # Выбираем layout
            layout = title_layout if idx == 0 else content_layout
            slide = prs.slides.add_slide(layout)

            # Заголовок — первая строка с # или просто первая
            title_text = ""
            body_lines: list[str] = []
            for line in lines:
                if line.startswith("# ") and not title_text:
                    title_text = line[2:].strip()
                elif line.startswith("## ") and not title_text:
                    title_text = line[3:].strip()
                else:
                    body_lines.append(line)

            if not title_text and lines:
                title_text = lines[0]
                body_lines = lines[1:]

            # Устанавливаем заголовок
            if slide.shapes.title:
                slide.shapes.title.text = title_text

            # Контент (только для слайдов с placeholder[1])
            if len(slide.placeholders) > 1 and body_lines:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                first = True
                for bl in body_lines:
                    if bl.startswith("- ") or bl.startswith("* ") or bl.startswith("• "):
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        p.text = bl[2:].strip()
                        p.level = 0
                    elif bl.startswith("  - ") or bl.startswith("  * "):
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        p.text = bl[4:].strip()
                        p.level = 1
                    else:
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        p.text = bl
                    first = False

        buf = io.BytesIO()
        prs.save(buf)
        return (
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    raise ValueError(f"Неизвестный формат: {ext}")
